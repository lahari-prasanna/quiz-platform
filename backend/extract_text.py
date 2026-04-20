import sys

def extract(file_path):
    ext = file_path.split('.')[-1].lower()
    text = ''
    try:
        if ext == 'pdf':
            import PyPDF2
            with open(file_path, 'rb') as f:
                reader = PyPDF2.PdfReader(f)
                for page in reader.pages:
                    text += page.extract_text() or ''
        elif ext == 'docx':
            from docx import Document
            doc = Document(file_path)
            for para in doc.paragraphs:
                text += para.text + '\n'
        elif ext in ['pptx', 'ppt']:
            from pptx import Presentation
            prs = Presentation(file_path)
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, 'text'):
                        text += shape.text + '\n'
        elif ext == 'txt':
            with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                text = f.read()
    except Exception as e:
        print(f'Error: {e}', file=sys.stderr)
    print(text)

if __name__ == '__main__':
    if len(sys.argv) < 2:
        sys.exit(1)
    extract(sys.argv[1])
